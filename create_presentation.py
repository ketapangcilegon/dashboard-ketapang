import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide dimensions to Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Constants
    COLOR_PRIMARY_DARK = RGBColor(0, 96, 56)      # Emerald Green #006038
    COLOR_PRIMARY_MED  = RGBColor(0, 122, 72)     # Emerald Green #007A48
    COLOR_ACCENT_GOLD   = RGBColor(217, 119, 6)    # Amber Gold #D97706
    COLOR_TEXT_DARK    = RGBColor(15, 23, 42)     # Slate Dark #0F172A
    COLOR_TEXT_MUTED   = RGBColor(71, 85, 105)    # Slate Muted #475569
    COLOR_BG_LIGHT     = RGBColor(248, 250, 252)  # Slate Light #F8FAFC
    COLOR_WHITE        = RGBColor(255, 255, 255)  # Pure White
    COLOR_CARD_BORDER  = RGBColor(226, 232, 240)  # Slate 200 #E2E8F0
    COLOR_GOLD_BG      = RGBColor(254, 249, 195)  # Pastel Gold #FEF9C3

    def add_header(slide, title_text, category_text="DASHBOARD KETAPANG KOTA CILEGON"):
        # Header background shape
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        header_shape.line.fill.background()

        # Category / Subtitle
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(167, 243, 208) # Mint Green

        # Main Title
        txBoxTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6))
        tfTitle = txBoxTitle.text_frame
        tfTitle.word_wrap = True
        pTitle = tfTitle.paragraphs[0]
        pTitle.text = title_text
        pTitle.font.size = Pt(22)
        pTitle.font.bold = True
        pTitle.font.color.rgb = COLOR_WHITE

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    bg1.line.fill.background()

    # Accent Gold Banner
    gold_banner = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.5), Inches(0.4))
    gold_banner.fill.solid()
    gold_banner.fill.fore_color.rgb = COLOR_ACCENT_GOLD
    gold_banner.line.fill.background()
    tf_gb = gold_banner.text_frame
    p_gb = tf_gb.paragraphs[0]
    p_gb.text = "SISTEM INFORMASI & EWS DIGITALE"
    p_gb.alignment = PP_ALIGN.CENTER
    p_gb.font.size = Pt(11)
    p_gb.font.bold = True
    p_gb.font.color.rgb = COLOR_WHITE

    # Title
    t_box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.2))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "KETAPANG KOTA CILEGON"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Sistem Peringatan Dini, Analisis 3 Pilar FSVA, EWS SKPG, & Peramalan Harga Pangan Berbasis Machine Learning"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(209, 250, 229)

    # Info Card Bottom
    info_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2))
    info_card.fill.solid()
    info_card.fill.fore_color.rgb = RGBColor(4, 77, 46)
    info_card.line.color.rgb = RGBColor(16, 185, 129)
    tf_ic = info_card.text_frame
    tf_ic.word_wrap = True
    p_ic1 = tf_ic.paragraphs[0]
    p_ic1.text = "Pemerintah Kota Cilegon | Dinas Ketahanan Pangan dan Pertanian"
    p_ic1.font.size = Pt(16)
    p_ic1.font.bold = True
    p_ic1.font.color.rgb = COLOR_WHITE

    p_ic2 = tf_ic.add_paragraph()
    p_ic2.text = "Dokumen Presentasi Eksekutif: Metodologi, Arsitektur Sistem, Integrasi Data FSVA 2025, Dampak Kebijakan, dan Referensi Ilmiah"
    p_ic2.font.size = Pt(13)
    p_ic2.font.color.rgb = RGBColor(167, 243, 208)

    # =========================================================================
    # SLIDE 2: GLOSARIUM & SINGKATAN UI/UX
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Daftar Singkatan & Glosarium Istilah UI/UX Sistem")

    table_shape = slide2.shapes.add_table(7, 4, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.4))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(4.1)

    headers = ["Singkatan", "Kepanjangan / Definisi", "Singkatan", "Kepanjangan / Definisi"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE

    acronyms = [
        ("FSVA", "Food Security and Vulnerability Atlas (Peta Ketahanan & Kerawanan Pangan)", "NCPR", "Net Cereals Per Capita Requirement (Rasio Pangan Pokok)"),
        ("IKP", "Indeks Ketahanan Pangan (Komposit Skala 0-100)", "AKE", "Angka Kecukupan Energi (Standar Ketersediaan Energi)"),
        ("SKPG", "Sistem Kewaspadaan Pangan dan Gizi (Pemantauan Bulanan)", "PoU", "Prevalence of Undernourishment (Kurang Konsumsi Energi)"),
        ("EWS", "Early Warning System (Sistem Peringatan Dini Kerawanan)", "PPH", "Pola Pangan Harapan (Keragaman & Kualitas Konsumsi)"),
        ("BPN / Bapanas", "Badan Pangan Nasional Republik Indonesia", "GPM", "Gerakan Pangan Murah (Operasi Pasar Intervensi Harga)"),
        ("SAGON", "Sistem Informasi Harga Bahan Pokok Pasar Tradisional Cilegon", "GBDT / OLS", "Gradient Boosted Decision Trees / Ordinary Least Squares")
    ]

    for row_idx, data in enumerate(acronyms, start=1):
        for col_idx in range(4):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE if row_idx % 2 == 1 else COLOR_BG_LIGHT
            p = cell.text_frame.paragraphs[0]
            p.text = data[col_idx]
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_DARK if col_idx % 2 == 0 else COLOR_TEXT_MUTED
            if col_idx % 2 == 0:
                p.font.bold = True

    # =========================================================================
    # SLIDE 3: LATAR BELAKANG & URGENSI
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Latar Belakang & Urgensi Pembangunan Platform KETAPANG")

    cards_data = [
        ("Karakteristik Kota Industri", "Kota Cilegon merupakan pusat industri manufaktur & bahan kimia dengan luas lahan pertanian yang relatif terbatas. Ketergantungan pasokan pangan luar daerah sangat tinggi (>85%).", COLOR_PRIMARY_MED),
        ("Kerentanan Gejolak Harga", "Fluktuasi harga bahan pangan pokok (cabai, bawang, beras, daging) akibat faktor cuaca ekstrem, disparitas pasar, dan biaya logistik memicu risiko inflasi daerah & penurunan daya beli masyarakat.", COLOR_ACCENT_GOLD),
        ("Kebutuhan Data Real-Time", "Sistem evaluasi manual berbasis kertas/spreadsheet memperlambat deteksi kerawanan pangan. Diperlukan platform digital terpadu berbasis AI & GIS untuk intervensi tepat waktu.", COLOR_PRIMARY_DARK)
    ]

    for idx, (title, desc, color) in enumerate(cards_data):
        left_pos = Inches(0.8 + idx * 3.9)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(3.7), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Header tag
        tag = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.5), Inches(3.7), Inches(0.8))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf_t = tag.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.alignment = PP_ALIGN.CENTER

        # Desc
        tx = slide3.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.5), Inches(3.3), Inches(4.1))
        tf_d = tx.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 4: RUMUSAN MASALAH & RISIKO JIKA TIDAK DIBUAT
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Masalah Utama yang Diselesaikan & Risiko Tanpa Web App")

    # Left Card: Masalah Utama
    card_l = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_WHITE
    card_l.line.color.rgb = COLOR_PRIMARY_DARK
    card_l.line.width = Pt(2)

    tx_l = slide4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_l = tx_l.text_frame
    tf_l.word_wrap = True
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "🎯 MASALAH UTAMA YANG DISLESAIKAN"
    p_lh.font.size = Pt(15)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_PRIMARY_DARK

    bullets_l = [
        "Fragmentasi Data Pangan: Data harga, gizi balita, dan iklim terpisah di berbagai instansi.",
        "Keterlambatan Deteksi Dini: Tidak ada instrumen otomatis yang mampu memprediksi gejolak harga 1-3 bulan ke depan.",
        "Ketidakakuratan Target Intervensi: Kesulitan menentukan kelurahan mana yang membutuhkan Bantuan Pangan / GPM secara mendesak.",
        "Kurangnya Akses Transparansi: Publik & pedagang kesulitan memantau tren harga resmi pasar."
    ]
    for b in bullets_l:
        p = tf_l.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_TEXT_DARK

    # Right Card: Risiko Tanpa Sistem
    card_r = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = RGBColor(254, 242, 242) # Pastel Red
    card_r.line.color.rgb = RGBColor(225, 29, 72) # Red 600
    card_r.line.width = Pt(2)

    tx_r = slide4.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_r = tx_r.text_frame
    tf_r.word_wrap = True
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "⚠️ RISIKO FATAL JIKA WEB APP TIDAK DIBUAT"
    p_rh.font.size = Pt(15)
    p_rh.font.bold = True
    p_rh.font.color.rgb = RGBColor(190, 18, 60)

    bullets_r = [
        "Lag Intervensi Kebijakan: Pemerintah daerah terlambat merespons lonjakan harga hingga menyebabkan panic buying.",
        "Eskalasi Angka Kerawanan Pangan: Kelurahan rentan tanpa air bersih & pendapatan rendah terlambat mendapat pasokan.",
        "Kerugian Ekonomi & Spikulasi Pasar: Pedagang dan konsumen terdampak fluktuasi harga tanpa kepastian stok.",
        "Kegagalan Pencapaian Target SPM: Ketidakmampuan memantau prevalensi stunting & kecukupan gizi balita secara berkala."
    ]
    for b in bullets_r:
        p = tf_r.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(159, 18, 57)

    # =========================================================================
    # SLIDE 5: FITUR-FITUR UTAMA & FUNGSI SISTEM
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Fitur Utama & Fungsi Platform Dashboard Ketapang")

    features = [
        ("🗺️ Peta Ketahanan Pangan (FSVA)", "Visualisasi GIS 43 Kelurahan berbasis IKP komposit & 3 pilar (Ketersediaan, Akses, Pemanfaatan).", COLOR_PRIMARY_DARK),
        ("🕸️ Grafik Radar Ketahanan Kelurahan", "Visualisasi jaring laba-laba 11 Indikator FSVA 2025 dengan mode komparasi Single & Dual Kelurahan.", COLOR_PRIMARY_MED),
        ("🤖 Peramalan Harga Pangan AI (ML)", "Prediksi harga 1 & 3 bulan ke depan berbasis Machine Learning (Random Forest, GBDT, OLS) + Fitur HBKN & Cuaca BMKG.", COLOR_ACCENT_GOLD),
        ("⚠️ Analisis SKPG & EWS Bulanan", "Sistem Peringatan Dini status kewaspadaan pangan bulanan kelurahan (Aman, Waspada, Rentan).", COLOR_PRIMARY_DARK),
        ("🛒 Scraping Harga Harian (SAGON)", "Integrasi data pencatatan harga bahan pokok harian dari pasar-pasar tradisional di Kota Cilegon.", COLOR_PRIMARY_MED),
        ("📋 Portal Input Admin & Rekomendasi", "Manajemen data upload FSVA, SKPG, Gizi Balita, serta penjanaan otomatis Rekomendasi Kebijakan (Policy Action).", COLOR_ACCENT_GOLD)
    ]

    for idx, (f_title, f_desc, f_color) in enumerate(features):
        row = idx // 3
        col = idx % 3
        l = Inches(0.8 + col * 3.9)
        t = Inches(1.5 + row * 2.7)

        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(3.7), Inches(2.4))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_WHITE
        c.line.color.rgb = f_color
        c.line.width = Pt(1.5)

        tx = slide5.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), Inches(3.4), Inches(2.1))
        tf = tx.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f_title
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = f_color

        p2 = tf.add_paragraph()
        p2.text = f_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 6: HIGHLIGHT GRAFIK RADAR 11 INDIKATOR FSVA 2025
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Highlight Fitur: Grafik Radar Ketahanan Pangan Kelurahan")

    # Left: 11 Indicators Breakdown
    box_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(6.0), Inches(5.4))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = COLOR_WHITE
    box_l.line.color.rgb = COLOR_PRIMARY_DARK
    box_l.line.width = Pt(2)

    tx_l = slide6.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.6), Inches(5.2))
    tf_l = tx_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "📌 11 INDIKATOR FSVA FORM 2 (JUKNIS BAPANAS 2025)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_DARK

    ind_text = [
        "1. Pilar Ketersediaan Pangan:",
        "   • 1.1 NCPR (Rasio Pangan Pokok) | 1.2 Ketersediaan Energi (AKE)",
        "   • 1.3 Protein Hewani | 1.4 Cadangan Pangan",
        "2. Pilar Keterjangkauan Pangan:",
        "   • 2.1 Penduduk Miskin (Desil 1+2) | 2.2 Stabilitas Harga (CV)",
        "   • 2.3 PoU (Prevalence of Undernourishment)",
        "3. Pilar Pemanfaatan Pangan:",
        "   • 3.1 Rata-rata Lama Sekolah Perempuan 15+ Thn",
        "   • 3.2 RT Tanpa Akses Air Bersih | 3.3 Skor Pola Pangan Harapan (PPH)",
        "   • 3.4 Prevalensi Balita Stunting"
    ]
    for line in ind_text:
        p_i = tf_l.add_paragraph()
        p_i.text = line
        p_i.font.size = Pt(10.5)
        p_i.font.color.rgb = COLOR_TEXT_DARK
        if line.startswith("1.") or line.startswith("2.") or line.startswith("3."):
            p_i.font.bold = True
            p_i.font.color.rgb = COLOR_PRIMARY_MED

    # Right: Pastle Gold Evaluation Box & Mode Switcher
    box_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.4), Inches(5.4), Inches(5.4))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = COLOR_GOLD_BG
    box_r.line.color.rgb = COLOR_ACCENT_GOLD
    box_r.line.width = Pt(2)

    tx_r = slide6.shapes.add_textbox(Inches(7.3), Inches(1.6), Inches(5.0), Inches(5.0))
    tf_r = tx_r.text_frame
    tf_r.word_wrap = True
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "✨ KARTU EVALUASI & REKOMENDASI KEBIJAKAN"
    p_rh.font.size = Pt(13)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_ACCENT_GOLD

    eval_items = [
        "• Single vs City Average Mode:",
        "  Membandingkan profil kelurahan terpilih secara otomatis terhadap Rata-rata Kota Cilegon.",
        "• Dual Kelurahan Comparison Mode:",
        "  Membandingkan 2 kelurahan secara head-to-head untuk penetapan prioritas intervensi.",
        "• Dynamic Policy Action (Rekomendasi Kebijakan):",
        "  Sistem secara otomatis mendeteksi indikator terlemah di kelurahan tersebut dan menghasilkan rekomendasi aksi spesifik (misal: PMT Lokal Posyandu untuk stunting, GPM untuk CV harga, atau perpipaan air bersih)."
    ]
    for ei in eval_items:
        p_e = tf_r.add_paragraph()
        p_e.text = ei
        p_e.font.size = Pt(11)
        p_e.font.color.rgb = COLOR_TEXT_DARK
        if ei.startswith("•"):
            p_e.font.bold = True

    # =========================================================================
    # SLIDE 7: TARGET AUDIENCE & PEMANFAAT
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Target Pengguna & Pemanfaat Utama Sistem")

    aud_data = [
        ("🏛️ Pemerintah Kota Cilegon & TPID", "Wali Kota, Sekda, dan Tim Pengendali Inflasi Daerah (TPID) menggunakan dashboard eksekutif sebagai landasan keputusan Kebijakan Ketahanan Pangan & Operasi Pasar Murah.", COLOR_PRIMARY_DARK),
        ("📊 Analis Ketahanan Pangan & OPD", "Analis Dinas Ketahanan Pangan (DKPP), Bappeda, dan Disperindag memanfaatkan rincian 11 indikator FSVA & EWS SKPG untuk perencanaan program intervensi spesifik.", COLOR_PRIMARY_MED),
        ("🎓 Akademisi & Peneliti", "Dosen, mahasiswa, dan peneliti menggunakan basis data historis 5 tahun, peramalan ML, serta indeks FSVA untuk studi ilmiah dan kebijakan publik.", COLOR_ACCENT_GOLD),
        ("👥 Masyarakat & Pelaku Usaha", "Warga Kota Cilegon dan pedagang pasar dapat secara transparan memantau perkembangan harga bahan pokok harian & kondisi ketahanan pangan wilayahnya.", COLOR_TEXT_DARK)
    ]

    for idx, (a_title, a_desc, a_color) in enumerate(aud_data):
        l = Inches(0.8 + (idx % 2) * 5.9)
        t = Inches(1.5 + (idx // 2) * 2.7)

        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(5.6), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = a_color
        card.line.width = Pt(2)

        tx = slide7.shapes.add_textbox(l + Inches(0.2), t + Inches(0.2), Inches(5.2), Inches(2.0))
        tf = tx.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = a_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = a_color

        p2 = tf.add_paragraph()
        p2.text = a_desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 8: ANALISIS DAMPAK MULTI-SEKTOR
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Analisis Dampak Strategis Multi-Sektor")

    table_shape8 = slide8.shapes.add_table(6, 2, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.4))
    t8 = table_shape8.table
    t8.columns[0].width = Inches(3.2)
    t8.columns[1].width = Inches(8.5)

    # Header
    for col_idx, htext in enumerate(["Entitas / Pemangku Kepentingan", "Dampak Positif & Nilai Tambah Platform KETAPANG"]):
        cell = t8.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = htext
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE

    impacts = [
        ("Ketahanan Pangan Kota Cilegon", "Meningkatnya ketahanan pangan daerah melalui deteksi dini kerawanan, stabilitas pasokan energi & protein, serta penurunan angka prevalensi stunting & PoU."),
        ("Pemkot Cilegon & TPID", "Mewujudkan tata kelola pemerintahan berbasis data (Data-Driven Governance), presisi dalam penganggaran APBD untuk bantuan pangan & GPM, serta efektivitas pengendalian inflasi daerah."),
        ("Analis Ketahanan Pangan", "Efisiensi waktu analisis dari hitungan minggu menjadi hitungan detik. Tersedianya instrumen otomatis pencetak rekomendasi kebijakan (Policy Action) per kelurahan."),
        ("Akademisi & Peneliti", "Tersedianya open dataset pangan daerah yang tervalidasi 5 tahun untuk riset ilmiah, pemodelan ekonomi pertanian, dan publikasi jurnal."),
        ("Masyarakat Umum", "Terlindunginya daya beli masyarakat dari gejolak harga yang tidak terkendali serta terjaminnya akses informasi pangan transparan.")
    ]

    for r_idx, (ent, imp) in enumerate(impacts, start=1):
        cell_e = t8.cell(r_idx, 0)
        cell_i = t8.cell(r_idx, 1)

        cell_e.fill.solid()
        cell_e.fill.fore_color.rgb = COLOR_BG_LIGHT if r_idx % 2 == 1 else COLOR_WHITE
        cell_i.fill.solid()
        cell_i.fill.fore_color.rgb = COLOR_BG_LIGHT if r_idx % 2 == 1 else COLOR_WHITE

        pe = cell_e.text_frame.paragraphs[0]
        pe.text = ent
        pe.font.bold = True
        pe.font.size = Pt(11.5)
        pe.font.color.rgb = COLOR_PRIMARY_DARK

        pi = cell_i.text_frame.paragraphs[0]
        pi.text = imp
        pi.font.size = Pt(11)
        pi.font.color.rgb = COLOR_TEXT_DARK

    # =========================================================================
    # SLIDE 9: METODOLOGI AI & ARSITEKTUR MESIN PREDIKSI
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Keandalan Mesin AI & Metodologi Peramalan Machine Learning")

    blocks = [
        ("🧠 Multi-Model Ensemble Engine", "Sistem membandingkan 3 algoritma secara independen untuk tiap komoditas:\n1. Random Forest Regressor\n2. Custom Gradient Boosting (GBDT)\n3. Multiple Linear Regression (OLS)\nChampion dipilih dari Walk-Forward MAPE terendah.", COLOR_PRIMARY_DARK),
        ("⚙️ Feature Engineering Lanjutan", "Mengintegrasikan 3 dimensi variabel:\n• Lag Harga Historis (t-1, t-2, t-3, MA3, MA6, Rolling Std)\n• Fitur Musiman HBKN (Ramadhan, Idul Fitri, Idul Adha, Nataru)\n• Iklim BMKG (Curah Hujan, Suhu, Kelembapan, Hari Hujan)", COLOR_PRIMARY_MED),
        ("🛡️ Serverless Native Execution & Log", "Berjalan 100% pada TypeScript Native Engine (Serverless Vercel) tanpa dependensi binary C++. Performa validasi tercatat otomatis di tabel Supabase ml_metrics (Akurasi Rata-rata 88.8% / MAPE 11.18%).", COLOR_ACCENT_GOLD)
    ]

    for idx, (b_title, b_desc, b_color) in enumerate(blocks):
        l = Inches(0.8 + idx * 3.9)
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, Inches(1.5), Inches(3.7), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = b_color
        card.line.width = Pt(2)

        tag = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.5), Inches(3.7), Inches(0.8))
        tag.fill.solid()
        tag.fill.fore_color.rgb = b_color
        tag.line.fill.background()
        tf_t = tag.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = b_title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.alignment = PP_ALIGN.CENTER

        tx = slide9.shapes.add_textbox(l + Inches(0.2), Inches(2.5), Inches(3.3), Inches(4.1))
        tf = tx.text_frame
        tf.word_wrap = True
        p_d = tf.paragraphs[0]
        p_d.text = b_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    # =========================================================================
    # SLIDE 10: KESIMPULAN & VAKSI STRATEGIS
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Kesimpulan & Arah Pengembangan Strategis")

    c_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    c_box.line.fill.background()

    tx = slide10.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.7))
    tf = tx.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "📌 KESIMPULAN UTAMA PLATFORM KETAPANG:"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE

    conclusions = [
        "1. Solusi Kompleks Kerawanan Pangan: Platform KETAPANG berhasil memodernisasi analisis ketahanan pangan Kota Cilegon melalui penggabungan Peta GIS 43 Kelurahan, Grafik Radar 11 Indikator FSVA 2025, dan EWS SKPG.",
        "2. Keandalan AI & Peringatan Dini: Penggunaan Machine Learning berbasis HBKN & Cuaca BMKG terbukti memberikan akurasi prediksi harga pangan sebesar 88.8% (MAPE 11.18%), memungkinkan Pemkot Cilegon melakukan aksi pencegahan sebelum gejolak harga terjadi.",
        "3. Efektivitas Pengambilan Kebijakan: Penjanaan otomatis Rekomendasi Kebijakan (Policy Action) memastikan setiap intervensi (GPM, PMT Lokal, Bantuan Pangan, Air Bersih) tepat sasaran di tingkat kelurahan.",
        "4. Menuju Smart Food City: Platform ini menempatkan Kota Cilegon sebagai pelopor tata kelola ketahanan pangan digital di Provinsi Banten."
    ]
    for c in conclusions:
        p = tf.add_paragraph()
        p.text = c
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(209, 250, 229)

    # =========================================================================
    # SLIDE 11: REFERENSI ILMIAH (APA STYLE 6TH EDITION)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Daftar Referensi Ilmiah (APA Style 6th Edition)")

    tx_ref = slide11.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.7))
    tf_ref = tx_ref.text_frame
    tf_ref.word_wrap = True

    references = [
        "Badan Pangan Nasional. (2025). Petunjuk teknis penyusunan Peta Ketahanan dan Kerawanan Pangan (Food Security and Vulnerability Atlas - FSVA) Kabupaten/Kota. Jakarta: Bapanas RI.",
        "Badan Pusat Statistik Kota Cilegon. (2025). Kota Cilegon dalam angka 2025. Cilegon: BPS Kota Cilegon.",
        "Food and Agriculture Organization. (2023). The state of food security and nutrition in the world 2023: Transforming food systems for affordable healthy diets. Rome: FAO.",
        "Hastuti, R., & Rahmanto, B. (2022). Penerapan Early Warning System (EWS) Sistem Kewaspadaan Pangan dan Gizi di tingkat daerah. Jurnal Analisis Kebijakan Pertanian, 20(1), 45-58.",
        "Hyndman, R. J., & Athanasopoulos, G. (2021). Forecasting: Principles and practice (3rd ed.). Melbourne: OTexts.",
        "Nitiyudo, A., & Purwanto, E. (2024). Pemodelan prediksi harga pangan pokok menggunakan algoritma Machine Learning berbasis faktor iklim dan musiman. Jurnal Teknologi Informasi dan Rekayasa Sistem, 12(2), 112-125.",
        "Sudaryanto, T., Syahyuti, S., & Agus, F. (2021). Strategi dan kebijakan ketahanan pangan wilayah perkotaan berbasis lahan terbatas. Jurnal Penelitian dan Pengembangan Pertanian, 40(2), 89-102.",
        "World Bank. (2023). Indonesia economic prospect: Climate change and food security. Washington, DC: World Bank Group."
    ]

    for idx, ref in enumerate(references):
        p = tf_ref.add_paragraph() if idx > 0 else tf_ref.paragraphs[0]
        p.text = ref
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_DARK

    # =========================================================================
    # SLIDE 12: CLOSING SLIDE
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    bg12.line.fill.background()

    t_box12 = slide12.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0))
    tf12 = t_box12.text_frame
    tf12.word_wrap = True
    p12_1 = tf12.paragraphs[0]
    p12_1.text = "TERIMA KASIH"
    p12_1.alignment = PP_ALIGN.CENTER
    p12_1.font.size = Pt(44)
    p12_1.font.bold = True
    p12_1.font.color.rgb = COLOR_WHITE

    p12_2 = tf12.add_paragraph()
    p12_2.text = "Platform KETAPANG Kota Cilegon - Mewujudkan Ketahanan Pangan Berkelanjutan & Berbasis Digital"
    p12_2.alignment = PP_ALIGN.CENTER
    p12_2.font.size = Pt(18)
    p12_2.font.color.rgb = RGBColor(209, 250, 229)

    p12_3 = tf12.add_paragraph()
    p12_3.text = "Dinas Ketahanan Pangan dan Pertanian Kota Cilegon"
    p12_3.alignment = PP_ALIGN.CENTER
    p12_3.font.size = Pt(14)
    p12_3.font.color.rgb = RGBColor(167, 243, 208)

    # Save PowerPoint file
    out_file = "Presentasi_Ketapang_Cilegon.pptx"
    prs.save(out_file)
    print(f"Successfully generated presentation deck: {out_file}")

if __name__ == '__main__':
    create_deck()
